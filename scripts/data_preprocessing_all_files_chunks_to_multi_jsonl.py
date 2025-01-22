# scripts/data_preprocessing.py

import os
import re
import pdfplumber
import docx
from bs4 import BeautifulSoup
import pytesseract
from PIL import Image
import json
import shutil
import zipfile
import tarfile
import csv
import io
from pptx import Presentation
import openpyxl
from striprtf.striprtf import rtf_to_text  # For RTF files

def clean_text(text: str) -> str:
    """
    Performs basic text cleaning:
    - Removes excessive whitespace
    - Cleans newline and other unwanted characters
    - Additional normalization if necessary
    """
    # Replace newline characters with space
    text = text.replace('\n', ' ')
    # Collapse multiple spaces into one
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    return text

def extract_text_from_pdf(pdf_path: str) -> str:
    """
    Extracts text from a PDF file.
    If the PDF is scanned (no text layer), OCR may be required.
    """
    full_text = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                cleaned = clean_text(text)
                full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_scanned_pdf(pdf_path: str) -> str:
    """
    If the PDF is scanned, converts each page to an image and uses OCR to extract text.
    This process is slower but necessary for scanned documents.
    """
    full_text = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            # Convert PDF page to PIL Image object
            pil_image = page.to_image(resolution=300).original
            # Extract text using Tesseract OCR
            ocr_text = pytesseract.image_to_string(pil_image, lang='tur')  # Ensure Turkish language model is available
            cleaned = clean_text(ocr_text)
            full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_docx(docx_path: str) -> str:
    """
    Reads text from a Word (.docx) file.
    """
    doc = docx.Document(docx_path)
    full_text = []
    for para in doc.paragraphs:
        cleaned = clean_text(para.text)
        if cleaned:
            full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_html(html_path: str) -> str:
    """
    Extracts and cleans the body text from an HTML file.
    """
    with open(html_path, 'r', encoding='utf-8') as f:
        html_content = f.read()
    soup = BeautifulSoup(html_content, 'lxml')
    # Extract only the body content
    text = soup.get_text(separator=' ')
    cleaned = clean_text(text)
    return cleaned

def extract_text_from_txt(txt_path: str) -> str:
    """
    Reads text from a .txt file.
    """
    with open(txt_path, 'r', encoding='utf-8') as f:
        text = f.read()
    cleaned = clean_text(text)
    return cleaned

def extract_text_from_rtf(rtf_path: str) -> str:
    """
    Reads text from a .rtf file.
    """
    with open(rtf_path, 'r', encoding='utf-8') as f:
        rtf_content = f.read()
    text = rtf_to_text(rtf_content)
    cleaned = clean_text(text)
    return cleaned

def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extracts text from a PowerPoint (.pptx) file.
    """
    prs = Presentation(pptx_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                cleaned = clean_text(shape.text)
                if cleaned:
                    full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_xlsx(xlsx_path: str) -> str:
    """
    Extracts text from an Excel (.xlsx) file.
    """
    wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)
    full_text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = ' '.join([str(cell) for cell in row if cell is not None])
            cleaned = clean_text(row_text)
            if cleaned:
                full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_csv(csv_path: str) -> str:
    """
    Reads text from a CSV file.
    """
    full_text = []
    with open(csv_path, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        for row in reader:
            row_text = ' '.join(row)
            cleaned = clean_text(row_text)
            if cleaned:
                full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_image(image_path: str) -> str:
    """
    Extracts text from image files (PNG, JPG, JPEG, TIFF, BMP) using OCR.
    """
    try:
        image = Image.open(image_path)
        ocr_text = pytesseract.image_to_string(image, lang='tur')  # Ensure Turkish language model is available
        cleaned = clean_text(ocr_text)
        return cleaned
    except Exception as e:
        print(f"Error processing image {image_path}: {e}")
        return ""

def chunk_text(text: str, chunk_size=500, overlap=50) -> list:
    """
    Splits text into chunks of approximately chunk_size words with a specified overlap.
    Overlapping words can improve retrieval quality.
    """
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk_words = words[start:end]
        chunk_str = " ".join(chunk_words)
        chunks.append(chunk_str)
        start += chunk_size - overlap
    return chunks

def save_chunks_to_jsonl(file_name: str, chunks: list, output_dir: str):
    """
    Saves text chunks to a JSON Lines (JSONL) file.
    Each line in the file is a separate JSON object.
    """
    out_path = os.path.join(output_dir, f"{file_name}.jsonl")
    with open(out_path, 'w', encoding='utf-8') as f:
        for i, chunk in enumerate(chunks):
            json_record = {
                "chunk_id": i,
                "text": chunk
            }
            f.write(json.dumps(json_record, ensure_ascii=False) + '\n')

def extract_text_from_zip(zip_path: str, temp_dir: str) -> list:
    """
    Extracts files from a ZIP archive and processes their text.
    """
    extracted_texts = []
    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                extracted_text = process_file(file_path, temp_dir)
                if extracted_text:
                    extracted_texts.append(extracted_text)
    except Exception as e:
        print(f"Error processing zip file {zip_path}: {e}")
    return extracted_texts

def extract_text_from_tar(tar_path: str, temp_dir: str) -> list:
    """
    Extracts files from a TAR archive and processes their text.
    """
    extracted_texts = []
    try:
        with tarfile.open(tar_path, 'r') as tar_ref:
            tar_ref.extractall(temp_dir)
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                extracted_text = process_file(file_path, temp_dir)
                if extracted_text:
                    extracted_texts.append(extracted_text)
    except Exception as e:
        print(f"Error processing tar file {tar_path}: {e}")
    return extracted_texts

def process_file(file_path: str, temp_dir: str) -> str:
    """
    Calls the appropriate text extraction function based on the file extension.
    """
    extension = os.path.splitext(file_path)[1].lower()
    try:
        if extension == ".pdf":
            text = extract_text_from_pdf(file_path)
            if not text.strip():
                text = extract_text_from_scanned_pdf(file_path)
            return text
        elif extension == ".docx":
            return extract_text_from_docx(file_path)
        elif extension == ".html":
            return extract_text_from_html(file_path)
        elif extension == ".txt":
            return extract_text_from_txt(file_path)
        elif extension == ".rtf":
            return extract_text_from_rtf(file_path)
        elif extension == ".pptx":
            return extract_text_from_pptx(file_path)
        elif extension == ".xlsx":
            return extract_text_from_xlsx(file_path)
        elif extension == ".csv":
            return extract_text_from_csv(file_path)
        elif extension in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
            return extract_text_from_image(file_path)
        elif extension == ".zip":
            return "\n".join(extract_text_from_zip(file_path, temp_dir))
        elif extension in [".tar", ".gz", ".bz2"]:
            return "\n".join(extract_text_from_tar(file_path, temp_dir))
        else:
            print(f"Unsupported file type: {extension}")
            return ""
    except Exception as e:
        print(f"Error processing file {file_path}: {e}")
        return ""

def preprocess_data_lake(data_lake_dir: str, output_dir: str):
    """
    Scans the data_lake_dir for various file types, extracts and cleans their text,
    and saves the processed data as JSON Lines (JSONL) files in the output_dir.
    """
    os.makedirs(output_dir, exist_ok=True)
    temp_dir = os.path.join(output_dir, "temp_extracted")
    os.makedirs(temp_dir, exist_ok=True)

    # Supported file extensions
    supported_extensions = {
        ".pdf", ".docx", ".html", ".txt", ".rtf",
        ".pptx", ".xlsx", ".csv",
        ".png", ".jpg", ".jpeg", ".tiff", ".bmp",
        ".zip", ".tar", ".gz", ".bz2"
    }

    for root, dirs, files in os.walk(data_lake_dir):
        for fname in files:
            extension = os.path.splitext(fname)[1].lower()
            if extension in supported_extensions:
                fpath = os.path.join(root, fname)
                print(f"Processing file: {fpath}")
                text_content = process_file(fpath, temp_dir)
                if text_content and text_content.strip():
                    chunks = chunk_text(text_content)
                    base_name = os.path.splitext(fname)[0]
                    save_chunks_to_jsonl(base_name, chunks, output_dir)

    # Clean up temporary directory
    shutil.rmtree(temp_dir)
    print(f"Data preprocessing completed. Cleaned JSONL files are stored in '{output_dir}' directory.")

if __name__ == "__main__":
    data_lake_directory = "../data_lake"
    cleaned_output_directory = "../cleaned_data"
    preprocess_data_lake(data_lake_directory, cleaned_output_directory)

"""
3.3 Attention to OCR
Some internal documents may be scanned PDFs. In this case, the extract_text_from_scanned_pdf function above should be used.
This process is slow and resource-intensive. Therefore, we may want to run it only on PDFs without a text layer.
"""
