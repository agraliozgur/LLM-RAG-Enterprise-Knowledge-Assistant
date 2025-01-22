# scripts/data_preprocessing.py

import os
import re
import pdfplumber
import docx
from bs4 import BeautifulSoup
import pytesseract
from PIL import Image
import json
import shutil

# Ek kütüphaneler
import zipfile
import tarfile
import csv
import io
from pptx import Presentation
import openpyxl
from striprtf.striprtf import rtf_to_text  # RTF dosyaları için

def clean_text(text: str) -> str:
    """
    Metin üzerinde basit temizlik yapar:
    - Fazla boşlukları giderir
    - Satır sonu vb. karakterleri temizler
    - Gerekiyorsa ek normalizasyon
    """
    # satır sonlarını boşlukla değiştir
    text = text.replace('\n', ' ')
    # birden fazla boşluk varsa teke indir
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    return text

def extract_text_from_pdf(pdf_path: str) -> str:
    """
    Bir PDF dosyasından metin çıkartan basit fonksiyon.
    Eğer PDF taranmışsa (metin katmanı yoksa) OCR gerekebilir.
    """
    full_text = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                cleaned = clean_text(text)
                full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_scanned_pdf(pdf_path: str) -> str:
    """
    Eğer PDF taranmışsa, pdfplumber page'lerinden her sayfayı resim olarak
    alıp OCR (tesseract) ile metin çekebiliriz.
    Bu işlem daha yavaştır ama taranmış belgeler için gereklidir.
    """
    full_text = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            # PDF sayfasını PIL Image objesine dönüştür
            pil_image = page.to_image(resolution=300).original
            # Tesseract OCR ile metin çıkart
            ocr_text = pytesseract.image_to_string(pil_image, lang='tur')  # Türkçe dil modeli varsa
            cleaned = clean_text(ocr_text)
            full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_docx(docx_path: str) -> str:
    """
    Word dosyalarından (docx) metin okur.
    """
    doc = docx.Document(docx_path)
    full_text = []
    for para in doc.paragraphs:
        cleaned = clean_text(para.text)
        if cleaned:
            full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_html(html_path: str) -> str:
    """
    HTML dosyasından body metnini temizleyerek alır.
    """
    with open(html_path, 'r', encoding='utf-8') as f:
        html_content = f.read()
    soup = BeautifulSoup(html_content, 'lxml')
    # sadece body kısımlarını çekelim
    text = soup.get_text(separator=' ')
    cleaned = clean_text(text)
    return cleaned

def extract_text_from_txt(txt_path: str) -> str:
    """
    .txt dosyalarından metin okur.
    """
    with open(txt_path, 'r', encoding='utf-8') as f:
        text = f.read()
    cleaned = clean_text(text)
    return cleaned

def extract_text_from_rtf(rtf_path: str) -> str:
    """
    .rtf dosyalarından metin okur.
    """
    with open(rtf_path, 'r', encoding='utf-8') as f:
        rtf_content = f.read()
    text = rtf_to_text(rtf_content)
    cleaned = clean_text(text)
    return cleaned

def extract_text_from_pptx(pptx_path: str) -> str:
    """
    PowerPoint dosyalarından (pptx) metin okur.
    """
    prs = Presentation(pptx_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                cleaned = clean_text(shape.text)
                if cleaned:
                    full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_xlsx(xlsx_path: str) -> str:
    """
    Excel dosyalarından (xlsx) metin okur.
    """
    wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)
    full_text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = ' '.join([str(cell) for cell in row if cell is not None])
            cleaned = clean_text(row_text)
            if cleaned:
                full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_csv(csv_path: str) -> str:
    """
    CSV dosyalarından metin okur.
    """
    full_text = []
    with open(csv_path, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        for row in reader:
            row_text = ' '.join(row)
            cleaned = clean_text(row_text)
            if cleaned:
                full_text.append(cleaned)
    return "\n".join(full_text)

def extract_text_from_image(image_path: str) -> str:
    """
    Görüntü dosyalarından (PNG, JPG, JPEG, TIFF, BMP) OCR ile metin çıkarır.
    """
    try:
        image = Image.open(image_path)
        ocr_text = pytesseract.image_to_string(image, lang='tur')  # Türkçe dil modeli
        cleaned = clean_text(ocr_text)
        return cleaned
    except Exception as e:
        print(f"Error processing image {image_path}: {e}")
        return ""

def chunk_text(text: str, chunk_size=500, overlap=50) -> list:
    """
    Metni chunk_size kelime civarında bölümler halinde parçalara ayırır.
    Overlap kadar kelime bir önceki chunk'la ortak bırakılır (retrieval kalitesini artırabilir).
    """
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk_words = words[start:end]
        chunk_str = " ".join(chunk_words)
        chunks.append(chunk_str)
        start += chunk_size - overlap
    return chunks

def save_chunks_to_json(file_name, chunks, output_dir):
    """
    Parçalanmış metinleri JSON formatında kaydeder.
    """
    data_to_save = []
    for i, ch in enumerate(chunks):
        data_to_save.append({
            "chunk_id": i,
            "text": ch
        })
    out_path = os.path.join(output_dir, file_name + ".json")
    with open(out_path, 'w', encoding='utf-8') as f:
        json.dump(data_to_save, f, ensure_ascii=False, indent=4)

def extract_text_from_zip(zip_path: str, temp_dir: str) -> list:
    """
    ZIP arşivinden dosyaları çıkarır ve metinlerini alır.
    """
    extracted_texts = []
    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                extracted_text = process_file(file_path, temp_dir)
                if extracted_text:
                    extracted_texts.append(extracted_text)
    except Exception as e:
        print(f"Error processing zip file {zip_path}: {e}")
    return extracted_texts

def extract_text_from_tar(tar_path: str, temp_dir: str) -> list:
    """
    TAR arşivinden dosyaları çıkarır ve metinlerini alır.
    """
    extracted_texts = []
    try:
        with tarfile.open(tar_path, 'r') as tar_ref:
            tar_ref.extractall(temp_dir)
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                extracted_text = process_file(file_path, temp_dir)
                if extracted_text:
                    extracted_texts.append(extracted_text)
    except Exception as e:
        print(f"Error processing tar file {tar_path}: {e}")
    return extracted_texts

def process_file(file_path: str, temp_dir: str) -> str:
    """
    Dosya türüne göre uygun metin çıkarma fonksiyonunu çağırır.
    """
    extension = os.path.splitext(file_path)[1].lower()
    try:
        if extension == ".pdf":
            text = extract_text_from_pdf(file_path)
            if not text.strip():
                text = extract_text_from_scanned_pdf(file_path)
            return text
        elif extension == ".docx":
            return extract_text_from_docx(file_path)
        elif extension == ".html":
            return extract_text_from_html(file_path)
        elif extension == ".txt":
            return extract_text_from_txt(file_path)
        elif extension == ".rtf":
            return extract_text_from_rtf(file_path)
        elif extension == ".pptx":
            return extract_text_from_pptx(file_path)
        elif extension == ".xlsx":
            return extract_text_from_xlsx(file_path)
        elif extension == ".csv":
            return extract_text_from_csv(file_path)
        elif extension in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
            return extract_text_from_image(file_path)
        elif extension in [".zip"]:
            return "\n".join(extract_text_from_zip(file_path, temp_dir))
        elif extension in [".tar", ".gz", ".bz2"]:
            return "\n".join(extract_text_from_tar(file_path, temp_dir))
        else:
            print(f"Unsupported file type: {extension}")
            return ""
    except Exception as e:
        print(f"Error processing file {file_path}: {e}")
        return ""

def preprocess_data_lake(data_lake_dir: str, output_dir: str):
    """
    data_lake_dir içinde çeşitli dosya türlerini tarar,
    metin olarak okuyup temizler ve output_dir'e JSON dosyaları olarak kaydeder.
    """
    os.makedirs(output_dir, exist_ok=True)
    temp_dir = os.path.join(output_dir, "temp_extracted")
    os.makedirs(temp_dir, exist_ok=True)

    # Desteklenen dosya türleri
    supported_extensions = {
        ".pdf", ".docx", ".html", ".txt", ".rtf",
        ".pptx", ".xlsx", ".csv",
        ".png", ".jpg", ".jpeg", ".tiff", ".bmp",
        ".zip", ".tar", ".gz", ".bz2"
    }

    for root, dirs, files in os.walk(data_lake_dir):
        for fname in files:
            extension = os.path.splitext(fname)[1].lower()
            if extension in supported_extensions:
                fpath = os.path.join(root, fname)
                print(f"Processing file: {fpath}")
                text_content = process_file(fpath, temp_dir)
                if text_content and text_content.strip():
                    chunks = chunk_text(text_content)
                    base_name = os.path.splitext(fname)[0]
                    save_chunks_to_json(base_name, chunks, output_dir)

    # Geçici dizini temizle
    shutil.rmtree(temp_dir)
    print(f"Veri ön işleme tamamlandı. Temizlenmiş JSON dosyaları '{output_dir}' klasöründe saklandı.")

if __name__ == "__main__":
    data_lake_directory = "../data_lake"
    cleaned_output_directory = "../cleaned_data"
    preprocess_data_lake(data_lake_directory, cleaned_output_directory)

"""
3.3 OCR’ye Dikkat
Kurum içi belgelerin bir kısmı taranmış PDF olabilir. Bu durumda yukarıdaki extract_text_from_scanned_pdf fonksiyonunu kullanmak gerekir.
Bu işlem yavaş ve kaynak tüketicidir. Dolayısıyla sadece metin katmanı boş olan PDF’lerde çalıştırmak isteyebiliriz.
"""
