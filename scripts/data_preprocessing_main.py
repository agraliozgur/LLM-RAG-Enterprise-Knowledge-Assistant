# scripts/data_preprocessing.py

import os
import re
import uuid
import json
import shutil
import zipfile
import tarfile
import csv
import pdfplumber
import docx
from pptx import Presentation
import openpyxl
from striprtf.striprtf import rtf_to_text  # For parsing RTF files
from bs4 import BeautifulSoup
import pytesseract
from PIL import Image
import nltk
from nltk.tokenize import sent_tokenize

# Ensure the required NLTK data is downloaded
nltk.download('punkt')

def clean_text(text: str) -> str:
    """
    Cleans the input text by removing URLs, HTML tags, unwanted characters,
    extra spaces, and leading/trailing whitespaces.

    Steps:
        1. Remove URLs (patterns like http://, https://, www.)
        2. Remove HTML tags
        3. Remove unwanted characters (keep alphanumeric and basic punctuation)
        4. Replace newlines with spaces
        5. Replace multiple spaces with a single space
        6. Trim leading and trailing spaces

    Args:
        text (str): The input text to be cleaned.

    Returns:
        str: The cleaned text.
    """
    if not isinstance(text, str):
        raise ValueError("Input text must be a string.")

    # 1. Remove URLs
    text = re.sub(r'http\S+|www\.\S+', '', text, flags=re.MULTILINE)

    # 2. Remove HTML tags
    text = re.sub(r'<.*?>', '', text)

    # 3. Remove unwanted characters (keep letters, digits, and punctuation)
    text = re.sub(r'[^A-Za-z0-9.,\'!?;:()\s]', '', text)

    # 4. Replace newline characters with space
    text = text.replace('\n', ' ')

    # 5. Replace multiple spaces with a single space
    text = re.sub(r'\s+', ' ', text)

    # 6. Strip leading and trailing spaces
    text = text.strip()

    return text

def extract_text_from_txt(txt_path: str) -> list:
    """
    Reads a .txt file line by line, cleans each valid line,
    and returns a list of chunks. A valid line is non-empty
    and has at least 100 characters.

    Args:
        txt_path (str): The path to the .txt file.

    Returns:
        list: A list of cleaned text lines (chunks).
    """
    chunks = []
    with open(txt_path, 'r', encoding='utf-8') as f:
        for line_number, line in enumerate(f, start=1):
            original_line = line.strip()
            if not original_line:
                print(f"Skipping empty line {line_number}.")
                continue
            if len(original_line) < 100:
                print(f"Skipping line {line_number} "
                      f"(length {len(original_line)} < 100).")
                continue
            cleaned_line = clean_text(original_line)
            if cleaned_line:
                chunks.append(cleaned_line)
            else:
                print(f"Line {line_number} was cleaned to an empty string. Skipping.")
    return chunks

def extract_text_from_pdf(pdf_path: str) -> str:
    """
    Extracts text from a PDF file using pdfplumber. If the PDF is
    scanned (no text layer), the result may be empty, in which case
    OCR might be required.

    Args:
        pdf_path (str): The path to the PDF file.

    Returns:
        str: The extracted and cleaned text, or an empty string on failure.
    """
    full_text = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    cleaned = clean_text(page_text)
                    full_text.append(cleaned)
    except Exception as e:
        print(f"Error processing PDF file {pdf_path}: {e}")

    return "\n".join(full_text)

def extract_text_from_scanned_pdf(pdf_path: str) -> str:
    """
    If the PDF is scanned (no text layer), converts each page to an image
    and applies OCR using Tesseract.

    Args:
        pdf_path (str): The path to the scanned PDF file.

    Returns:
        str: The extracted and cleaned text from all pages.
    """
    full_text = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                try:
                    # Convert the PDF page to a PIL Image
                    pil_image = page.to_image(resolution=300).original
                    # Apply OCR using Tesseract
                    ocr_text = pytesseract.image_to_string(pil_image, lang='eng')
                    cleaned = clean_text(ocr_text)
                    if cleaned:
                        full_text.append(cleaned)
                except Exception as inner_e:
                    print(f"Error performing OCR on page {page_number} "
                          f"of {pdf_path}: {inner_e}")
    except Exception as e:
        print(f"Error processing scanned PDF file {pdf_path}: {e}")

    return "\n".join(full_text)

def extract_text_from_docx(docx_path: str) -> str:
    """
    Extracts text from a .docx file using python-docx.

    Args:
        docx_path (str): The path to the .docx file.

    Returns:
        str: The extracted and cleaned text, or an empty string on failure.
    """
    try:
        document = docx.Document(docx_path)
        full_text = []
        for para in document.paragraphs:
            cleaned = clean_text(para.text)
            if cleaned:
                full_text.append(cleaned)
        return "\n".join(full_text)
    except Exception as e:
        print(f"Error processing DOCX file {docx_path}: {e}")
        return ""

def extract_text_from_html(html_path: str) -> str:
    """
    Extracts text from an HTML file, focusing on body content.

    Args:
        html_path (str): The path to the .html file.

    Returns:
        str: The extracted and cleaned text, or an empty string on failure.
    """
    try:
        with open(html_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        soup = BeautifulSoup(html_content, 'lxml')
        # Get all text with a separator of spaces
        text = soup.get_text(separator=' ')
        cleaned = clean_text(text)
        return cleaned
    except Exception as e:
        print(f"Error processing HTML file {html_path}: {e}")
        return ""

def extract_text_from_rtf(rtf_path: str) -> str:
    """
    Extracts text from a .rtf file using the striprtf library.

    Args:
        rtf_path (str): The path to the .rtf file.

    Returns:
        str: The extracted and cleaned text, or an empty string on failure.
    """
    try:
        with open(rtf_path, 'r', encoding='utf-8') as f:
            rtf_content = f.read()
        text = rtf_to_text(rtf_content)
        cleaned = clean_text(text)
        return cleaned
    except Exception as e:
        print(f"Error processing RTF file {rtf_path}: {e}")
        return ""

def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extracts text from a .pptx file using python-pptx.

    Args:
        pptx_path (str): The path to the .pptx file.

    Returns:
        str: The extracted and cleaned text, or an empty string on failure.
    """
    try:
        presentation = Presentation(pptx_path)
        full_text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    cleaned = clean_text(shape.text)
                    if cleaned:
                        full_text.append(cleaned)
        return "\n".join(full_text)
    except Exception as e:
        print(f"Error processing PPTX file {pptx_path}: {e}")
        return ""

def extract_text_from_xlsx(xlsx_path: str) -> str:
    """
    Extracts text from an Excel (.xlsx) file using openpyxl.

    Args:
        xlsx_path (str): The path to the .xlsx file.

    Returns:
        str: The extracted and cleaned text, or an empty string on failure.
    """
    try:
        workbook = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)
        full_text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join(str(cell) for cell in row if cell is not None)
                cleaned = clean_text(row_text)
                if cleaned:
                    full_text.append(cleaned)
        return "\n".join(full_text)
    except Exception as e:
        print(f"Error processing XLSX file {xlsx_path}: {e}")
        return ""

def extract_text_from_csv(csv_path: str) -> str:
    """
    Extracts text from a CSV file. Each row is joined into a single line.

    Args:
        csv_path (str): The path to the .csv file.

    Returns:
        str: The extracted and cleaned text, or an empty string on failure.
    """
    try:
        full_text = []
        with open(csv_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                row_text = ' '.join(row)
                cleaned = clean_text(row_text)
                if cleaned:
                    full_text.append(cleaned)
        return "\n".join(full_text)
    except Exception as e:
        print(f"Error processing CSV file {csv_path}: {e}")
        return ""

def extract_text_from_image(image_path: str) -> str:
    """
    Extracts text from image files (PNG, JPG, JPEG, TIFF, BMP) using OCR (Tesseract).

    Args:
        image_path (str): The path to the image file.

    Returns:
        str: The extracted and cleaned text, or an empty string on failure.
    """
    try:
        image = Image.open(image_path)
        ocr_text = pytesseract.image_to_string(image, lang='eng')
        cleaned = clean_text(ocr_text)
        return cleaned
    except Exception as e:
        print(f"Error processing image {image_path}: {e}")
        return ""

def process_file(file_path: str, temp_dir: str) -> str:
    """
    Chooses the appropriate extraction function based on file extension,
    then returns the extracted and cleaned text as a single string.

    Args:
        file_path (str): Path to the file to be processed.
        temp_dir (str): Temporary directory for extracted files (archives).

    Returns:
        str: Extracted and cleaned text.
    """
    extension = os.path.splitext(file_path)[1].lower()
    try:
        if extension == ".pdf":
            # Attempt direct PDF text extraction
            text = extract_text_from_pdf(file_path)
            if not text.strip():
                # If empty, attempt OCR on a scanned PDF
                text = extract_text_from_scanned_pdf(file_path)
            return text
        elif extension == ".docx":
            return extract_text_from_docx(file_path)
        elif extension == ".html":
            return extract_text_from_html(file_path)
        elif extension == ".txt":
            # .txt returns a list of chunks, so we'll handle it differently below
            txt_chunks = extract_text_from_txt(file_path)
            return "\n".join(txt_chunks)
        elif extension == ".rtf":
            return extract_text_from_rtf(file_path)
        elif extension == ".pptx":
            return extract_text_from_pptx(file_path)
        elif extension == ".xlsx":
            return extract_text_from_xlsx(file_path)
        elif extension == ".csv":
            return extract_text_from_csv(file_path)
        elif extension in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
            return extract_text_from_image(file_path)
        elif extension == ".zip":
            return "\n".join(extract_text_from_zip(file_path, temp_dir))
        elif extension in [".tar", ".gz", ".bz2"]:
            return "\n".join(extract_text_from_tar(file_path, temp_dir))
        else:
            print(f"Unsupported file type: {extension}")
            return ""
    except Exception as e:
        print(f"Error processing file {file_path}: {e}")
        return ""

def extract_text_from_zip(zip_path: str, temp_dir: str) -> list:
    """
    Extracts files from a ZIP archive to a temporary directory,
    processes each file, and returns the combined extracted text
    as a list of strings.

    Args:
        zip_path (str): Path to the .zip file.
        temp_dir (str): Temporary directory for extraction.

    Returns:
        list: A list of extracted text from the archive.
    """
    texts = []
    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                extracted_text = process_file(file_path, temp_dir)
                if extracted_text:
                    texts.append(extracted_text)
    except Exception as e:
        print(f"Error processing ZIP file {zip_path}: {e}")
    return texts

def extract_text_from_tar(tar_path: str, temp_dir: str) -> list:
    """
    Extracts files from a TAR/TAR.GZ/TAR.BZ2 archive to a temporary directory,
    processes each file, and returns the combined extracted text as a list of strings.

    Args:
        tar_path (str): Path to the tar/tar.gz/tar.bz2 file.
        temp_dir (str): Temporary directory for extraction.

    Returns:
        list: A list of extracted text from the archive.
    """
    texts = []
    try:
        with tarfile.open(tar_path, 'r') as tar_ref:
            tar_ref.extractall(temp_dir)
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                extracted_text = process_file(file_path, temp_dir)
                if extracted_text:
                    texts.append(extracted_text)
    except Exception as e:
        print(f"Error processing TAR file {tar_path}: {e}")
    return texts

def chunk_text_sentence_based_adjusted(
    text: str, 
    chunk_size: int = 200, 
    min_extra: int = 100
) -> list:
    """
    Splits text into chunks of approximately 'chunk_size' words, ensuring no chunk 
    is shorter than 'min_extra' words unless it's the only chunk.

    - Uses sentence tokenization (NLTK) to preserve sentence boundaries.
    - Merges the last two chunks if the final chunk is shorter than 'min_extra'.

    Args:
        text (str): The text to be split into chunks.
        chunk_size (int): The maximum number of words in each chunk.
        min_extra (int): The minimum number of words allowed in the last chunk.

    Returns:
        list: A list of text chunks that respect sentence boundaries.
    """
    # Tokenize sentences in English
    sentences = sent_tokenize(text, language='english')
    chunks = []
    current_chunk = []
    current_length = 0

    for sentence in sentences:
        sentence_words = sentence.split()
        sentence_length = len(sentence_words)

        # If a single sentence exceeds chunk_size, treat it as a separate chunk
        if sentence_length > chunk_size:
            # Finalize the current chunk if it has content
            if current_chunk:
                chunks.append(' '.join(current_chunk))
                current_chunk = []
                current_length = 0
            chunks.append(sentence)
            continue

        # If adding this sentence would exceed the chunk_size, finalize the current chunk
        if current_length + sentence_length > chunk_size:
            if current_chunk:
                chunks.append(' '.join(current_chunk))
            current_chunk = [sentence]
            current_length = sentence_length
        else:
            current_chunk.append(sentence)
            current_length += sentence_length

    # Add the final chunk if present
    if current_chunk:
        chunks.append(' '.join(current_chunk))

    # Check if the last chunk is too short
    if len(chunks) >= 2:
        last_chunk_word_count = len(chunks[-1].split())
        if last_chunk_word_count < min_extra:
            # Merge the last chunk with the previous one
            previous_chunk = chunks[-2]
            last_chunk = chunks[-1]
            merged_length = len(previous_chunk.split()) + last_chunk_word_count

            # If merging doesn't exceed chunk_size, merge them
            if merged_length <= chunk_size:
                chunks[-2] = previous_chunk + ' ' + last_chunk
                chunks.pop()
            else:
                # Otherwise, still merge if needed, or handle differently
                chunks[-2] = previous_chunk + ' ' + last_chunk
                chunks.pop()

    return chunks

def save_chunks_to_jsonl(
    unique_id: str, 
    file_name: str, 
    extension: str, 
    chunks: list, 
    output_file
):
    """
    Saves text chunks to a single JSON Lines (JSONL) file.
    Each chunk is written as a separate JSON object with metadata.

    Args:
        unique_id (str): A unique identifier for the file.
        file_name (str): Original file name.
        extension (str): File extension (e.g., '.txt', '.pdf').
        chunks (list): List of text chunks to be saved.
        output_file: A file object to write JSONL entries into.
    """
    for i, chunk in enumerate(chunks):
        json_record = {
            "unique_id": unique_id,
            "file_name": file_name,
            "extension": extension,
            "chunk_id": i,
            "text": chunk
        }
        output_file.write(json.dumps(json_record, ensure_ascii=False) + '\n')

def load_processed_files(log_file_path: str) -> set:
    """
    Loads a set of already processed file paths from the specified log file.

    Args:
        log_file_path (str): The path to the log file.

    Returns:
        set: A set of file paths that have already been processed.
    """
    processed_files = set()
    if os.path.exists(log_file_path):
        with open(log_file_path, 'r', encoding='utf-8') as log_file:
            for line in log_file:
                processed_file = line.strip()
                if processed_file:
                    processed_files.add(processed_file)
    return processed_files

def append_processed_file(log_file_path: str, file_path: str):
    """
    Appends a processed file path to the log file, indicating it has been handled.

    Args:
        log_file_path (str): Path to the log file.
        file_path (str): Path to the file that was processed.
    """
    with open(log_file_path, 'a', encoding='utf-8') as log_file:
        log_file.write(file_path + '\n')

def preprocess_data_lake(
    data_lake_dir: str, 
    output_file_path: str, 
    log_file_path: str
):
    """
    Traverses the specified data lake directory, extracts and cleans text 
    from supported file types, chunks the text, and saves it as JSON Lines.

    Also tracks processed files to avoid duplicate work across multiple runs.

    Args:
        data_lake_dir (str): Path to the data lake directory.
        output_file_path (str): Path to the output JSONL file.
        log_file_path (str): Path to the log file for tracking processed files.
    """
    # Ensure the output directory exists
    os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
    temp_dir = os.path.join(os.path.dirname(output_file_path), "temp_extracted")
    os.makedirs(temp_dir, exist_ok=True)

    # Supported file extensions
    supported_extensions = {
        ".pdf", ".docx", ".html", ".txt", ".rtf",
        ".pptx", ".xlsx", ".csv",
        ".png", ".jpg", ".jpeg", ".tiff", ".bmp",
        ".zip", ".tar", ".gz", ".bz2"
    }

    # Load already processed files to skip duplicates
    processed_files = load_processed_files(log_file_path)
    print(f"Loaded {len(processed_files)} already processed files from log: {log_file_path}")

    with open(output_file_path, 'a', encoding='utf-8') as output_file:
        for root, dirs, files in os.walk(data_lake_dir):
            for fname in files:
                extension = os.path.splitext(fname)[1].lower()
                if extension in supported_extensions:
                    fpath = os.path.abspath(os.path.join(root, fname))
                    if fpath in processed_files:
                        print(f"Skipping already processed file: {fpath}")
                        continue

                    print(f"Processing file: {fpath}")
                    unique_id = str(uuid.uuid4())
                    text_content = process_file(fpath, temp_dir)

                    # If text content is valid
                    if text_content and text_content.strip():
                        # For .txt files, each line is considered a chunk
                        if extension == ".txt":
                            txt_chunks = extract_text_from_txt(fpath)
                            chunks = txt_chunks
                        else:
                            # For other file types, use sentence-based chunking
                            chunks = chunk_text_sentence_based_adjusted(text_content)

                        if chunks:
                            save_chunks_to_jsonl(unique_id, fname, extension, chunks, output_file)
                            append_processed_file(log_file_path, fpath)
                            print(f"Logged processed file: {fpath}")
                    else:
                        print(f"No valid text extracted from: {fpath}")

    # Remove the temporary extraction directory
    shutil.rmtree(temp_dir)
    print(f"Data preprocessing completed. All chunks saved to: {output_file_path}")
    print(f"Processed files are recorded in: {log_file_path}")

if __name__ == "__main__":
    # Example usage: adjust paths as needed
    data_lake_directory = "../data/raw"  # Path to your data lake directory
    output_jsonl_file = "../data/cleaned_data/all_processed_data.jsonl"  # Output JSONL file path
    processed_files_log = "../data/cleaned_data/all_processed_files.log"  # Log file path

    preprocess_data_lake(data_lake_directory, output_jsonl_file, processed_files_log)
